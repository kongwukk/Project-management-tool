from __future__ import annotations

import csv
from io import StringIO
from pathlib import Path
from typing import Iterable, List

from .models import ParsedDocument


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx"}


class FileProcessingError(RuntimeError):
    pass


class UnsupportedFileType(FileProcessingError):
    pass


class MissingDependencyError(FileProcessingError):
    pass


def parse_office_file(path: str | Path) -> ParsedDocument:
    source = Path(path)
    suffix = source.suffix.lower()
    if suffix == ".docx":
        return ParsedDocument(source.name, "Word", _parse_docx(source), str(source))
    if suffix == ".pptx":
        return ParsedDocument(source.name, "PowerPoint", _parse_pptx(source), str(source))
    if suffix == ".xlsx":
        return ParsedDocument(source.name, "Excel", _parse_xlsx(source), str(source))
    raise UnsupportedFileType(f"不支持的文件格式：{suffix or '未知'}。请上传 .docx、.pptx 或 .xlsx 文件。")


def parse_office_files(paths: Iterable[str | Path]) -> List[ParsedDocument]:
    return [parse_office_file(path) for path in paths]


def _parse_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise MissingDependencyError("缺少 python-docx 依赖，请先运行：pip install python-docx") from exc

    document = Document(str(path))
    sections: List[str] = [f"# 文件：{path.name}", "## 段落"]

    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    sections.extend(paragraphs or ["（未提取到段落文本）"])

    if document.tables:
        sections.append("## 表格")
        for table_index, table in enumerate(document.tables, start=1):
            sections.append(f"### 表格 {table_index}")
            for row in table.rows:
                cells = [_clean_cell(cell.text) for cell in row.cells]
                sections.append(_to_markdown_row(cells))

    return "\n".join(sections)


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise MissingDependencyError("缺少 python-pptx 依赖，请先运行：pip install python-pptx") from exc

    presentation = Presentation(str(path))
    sections: List[str] = [f"# 文件：{path.name}"]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        sections.append(f"## 幻灯片 {slide_index}")
        texts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        sections.extend(texts or ["（未提取到幻灯片文本）"])

        notes = _extract_slide_notes(slide)
        if notes:
            sections.append("### 备注")
            sections.append(notes)

    if len(sections) == 1:
        sections.append("（未提取到幻灯片内容）")
    return "\n".join(sections)


def _parse_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise MissingDependencyError("缺少 openpyxl 依赖，请先运行：pip install openpyxl") from exc

    workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    sections: List[str] = [f"# 文件：{path.name}"]

    for sheet in workbook.worksheets:
        sections.append(f"## 工作表：{sheet.title}")
        rows_written = 0
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if not any(value.strip() for value in values):
                continue
            sections.append(_to_csv_row(values))
            rows_written += 1
        if rows_written == 0:
            sections.append("（空工作表）")

    workbook.close()
    return "\n".join(sections)


def _extract_slide_notes(slide: object) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
        return notes_frame.text.strip()
    except Exception:
        return ""


def _clean_cell(value: str) -> str:
    return " ".join(value.split())


def _to_markdown_row(values: List[str]) -> str:
    escaped = [value.replace("|", "\\|") for value in values]
    return "| " + " | ".join(escaped) + " |"


def _to_csv_row(values: List[str]) -> str:
    stream = StringIO()
    writer = csv.writer(stream)
    writer.writerow(values)
    return stream.getvalue().strip("\r\n")

